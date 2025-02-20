import os
from pathlib import Path
import fitz  # PyMuPDF for PDFs
from pptx import Presentation
import speech_recognition as sr
from moviepy.editor import AudioFileClip
from qdrant_client import QdrantClient
from qdrant_client.http import models
from sentence_transformers import SentenceTransformer
from config import Config

class FileProcessor:
    def __init__(self):
        self.client = QdrantClient(Config.QDRANT_HOST, port=Config.QDRANT_PORT)
        self.collection_name = Config.COLLECTION_NAME
        self.model = SentenceTransformer("all-MiniLM-L6-v2")

        # Ensure collection exists
        self.client.recreate_collection(
            collection_name=self.collection_name,
            vectors_config=models.VectorParams(size=384, distance=models.Distance.COSINE)
        )

    def process_pdf(self, file_path: str) -> str:
        """Extract text from a PDF file"""
        text = ""
        with fitz.open(file_path) as doc:
            for page in doc:
                text += page.get_text()
        return text

    def process_ppt(self, file_path: str) -> str:
        """Extract text from PowerPoint slides"""
        text = ""
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    def process_audio(self, file_path: str) -> str:
        """Convert audio to text"""
        recognizer = sr.Recognizer()
        with sr.AudioFile(file_path) as source:
            audio = recognizer.record(source)
            try:
                return recognizer.recognize_google(audio)
            except (sr.UnknownValueError, sr.RequestError):
                return "[Audio not recognized]"

    def process_video(self, file_path: str) -> str:
        """Extract audio from video and convert to text"""
        audio_path = "temp_audio.wav"
        video = AudioFileClip(file_path)
        video.write_audiofile(audio_path, codec="pcm_s16le")
        
        text = self.process_audio(audio_path)
        os.remove(audio_path)
        return text

    def store_file(self, file_path: str) -> bool:
        """Process and store file in Qdrant"""
        file_path = Path(file_path)
        
        if file_path.suffix.lower() == ".pdf":
            text = self.process_pdf(str(file_path))
        elif file_path.suffix.lower() in [".ppt", ".pptx"]:
            text = self.process_ppt(str(file_path))
        elif file_path.suffix.lower() in [".mp3", ".wav"]:
            text = self.process_audio(str(file_path))
        elif file_path.suffix.lower() in [".mp4", ".avi"]:
            text = self.process_video(str(file_path))
        else:
            raise ValueError(f"Unsupported file type: {file_path.suffix}")

        # Generate embedding
        embedding = self.model.encode(text).tolist()

        # Store in Qdrant
        self.client.upsert(
            collection_name=self.collection_name,
            points=[models.PointStruct(
                id=abs(hash(str(file_path))) % (10**12),
                vector=embedding,
                payload={"text": text, "filename": file_path.name, "file_type": file_path.suffix}
            )]
        )
        return True

    def search_similar_files(self, query: str, limit: int = 5):
        """Search for similar files based on text query"""
        query_vector = self.model.encode(query).tolist()
        
        search_result = self.client.search(
            collection_name=self.collection_name,
            query_vector=query_vector,
            with_payload=True,
            limit=limit
        )
        
        return [{"score": hit.score, "metadata": hit.payload} for hit in search_result]

    def get_text_from_qdrant(self, query: str, limit: int = 3):
        """Retrieve text content from Qdrant using similarity search"""
        search_results = self.search_similar_files(query, limit)

        if not search_results:
            return "No relevant content found in Qdrant."

        return " ".join([result["metadata"].get("text", "") for result in search_results if result["metadata"].get("text")])
